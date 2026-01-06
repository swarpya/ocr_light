from fastapi import FastAPI, File, UploadFile, Form
from fastapi.responses import JSONResponse
import io
import gc
import logging
from contextlib import asynccontextmanager
from PIL import Image
import pytesseract
import pypdfium2 as pdfium
import pandas as pd
from pptx import Presentation

# Configure Logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

@asynccontextmanager
async def lifespan(app: FastAPI):
    logger.info("--- STARTUP: ECONOMY MODE (TESSERACT ONLY) ---")
    logger.info("--- READY TO PROCESS ON 1GB RAM ---")
    yield
    logger.info("--- SHUTDOWN ---")

app = FastAPI(lifespan=lifespan)

@app.post("/ocr", response_class=JSONResponse)
async def process_document(
    file: UploadFile = File(...)
):
    content = await file.read()
    filename = file.filename.lower()
    
    # Standard Response Structure
    response_data = {
        "filename": filename, 
        "engine": "tesseract-economy", 
        "pages": []
    }

    try:
        # --- HANDLER 1: EXCEL (.xlsx, .xls) ---
        # Very RAM efficient, reads data directly.
        if filename.endswith(('.xlsx', '.xls')):
            logger.info(f"Processing Excel: {filename}")
            excel_file = io.BytesIO(content)
            xls = pd.ExcelFile(excel_file)
            
            for i, sheet_name in enumerate(xls.sheet_names):
                # Read sheet into dataframe
                df = pd.read_excel(xls, sheet_name=sheet_name)
                # Convert to string representation for content
                text_content = df.to_string(index=False)
                # Convert to raw data for JSON
                raw_json = df.fillna("").to_dict(orient='records')
                
                response_data["pages"].append({
                    "page": i+1,
                    "sheet_name": sheet_name,
                    "elements": [{
                        "type": "Table",
                        "content": text_content,
                        "data": raw_json
                    }]
                })
            return JSONResponse(content=response_data)

        # --- HANDLER 2: POWERPOINT (.pptx) ---
        # Extracts text from slides directly.
        elif filename.endswith('.pptx'):
            logger.info(f"Processing PowerPoint: {filename}")
            ppt_file = io.BytesIO(content)
            prs = Presentation(ppt_file)
            
            for i, slide in enumerate(prs.slides):
                page_elements = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        # Simple Heuristic for Labeling
                        elem_type = "Text"
                        if "title" in shape.name.lower():
                            elem_type = "Title"
                        
                        page_elements.append({
                            "type": elem_type,
                            "content": shape.text.strip()
                        })
                
                response_data["pages"].append({
                    "page": i+1,
                    "elements": page_elements
                })
            return JSONResponse(content=response_data)

        # --- HANDLER 3: PDF / IMAGES (TESSERACT OCR) ---
        else:
            logger.info(f"Processing Image/PDF: {filename}")
            pil_images = []
            
            # Load PDF or Image
            if filename.endswith(".pdf"):
                pdf = pdfium.PdfDocument(content)
                for i in range(len(pdf)):
                    # SCALE 1.5: Good balance of speed vs accuracy for Tesseract
                    pil_images.append(pdf[i].render(scale=1.5).to_pil().convert("RGB"))
            else:
                pil_images.append(Image.open(io.BytesIO(content)).convert("RGB"))

            # Process pages SEQUENTIALLY to save RAM
            for i, image in enumerate(pil_images):
                logger.info(f"   -> OCR Scanning Page {i+1}...")
                
                # Tesseract OCR
                # timeout=20 prevents server from hanging on bad images
                text = pytesseract.image_to_string(image, timeout=20)
                
                # Basic Post-Processing to simulate "Elements"
                # We split by double newlines to find paragraphs
                raw_blocks = text.split("\n\n")
                elements = []
                
                for block in raw_blocks:
                    clean_block = block.strip()
                    if clean_block:
                        # Simple Heuristic: Short + No Lowercase = Title?
                        # Otherwise default to "Text"
                        elem_type = "Text"
                        if len(clean_block) < 50 and clean_block.isupper():
                            elem_type = "Title"
                            
                        elements.append({
                            "type": elem_type,
                            "content": clean_block
                        })

                # Fallback if page is empty
                if not elements:
                    elements.append({"type": "Empty-Page", "content": ""})

                response_data["pages"].append({
                    "page": i+1,
                    "elements": elements
                })
                
                # AGGRESSIVE MEMORY CLEANUP
                del image
                gc.collect()

    except Exception as e:
        logger.error(f"Error processing file: {e}")
        return JSONResponse(content={"error": str(e)}, status_code=500)

    return JSONResponse(content=response_data)